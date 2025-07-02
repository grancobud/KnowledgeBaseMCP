"""
Document text extraction utilities for various file formats
"""

import os
import logging
from pathlib import Path
from typing import Dict, List, Optional, Union
import asyncio
from datetime import datetime

# Document processing imports
try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import pdfplumber
    PDF_AVAILABLE = True
    PDF_LIBRARY = 'pdfplumber'
except ImportError:
    try:
        import PyPDF2
        PDF_AVAILABLE = True
        PDF_LIBRARY = 'pypdf2'
    except ImportError:
        PDF_AVAILABLE = False
        PDF_LIBRARY = None

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import pandas as pd
    from openpyxl import load_workbook
    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False

logger = logging.getLogger(__name__)

class DocumentExtractor:
    """Main class for extracting text from various document formats"""
    
    SUPPORTED_EXTENSIONS = {
        '.pdf': 'PDF Document',
        '.docx': 'Word Document', 
        '.pptx': 'PowerPoint Presentation',
        '.xlsx': 'Excel Spreadsheet'
    }
    
    def __init__(self):
        self.check_dependencies()
    
    def check_dependencies(self):
        """Check if required libraries are available"""
        missing = []
        if not DOCX_AVAILABLE:
            missing.append("python-docx (for DOCX files)")
        if not PDF_AVAILABLE:
            missing.append("pdfplumber or PyPDF2 (for PDF files)")
        if not PPTX_AVAILABLE:
            missing.append("python-pptx (for PPTX files)")
        if not XLSX_AVAILABLE:
            missing.append("openpyxl and pandas (for XLSX files)")
        
        if missing:
            logger.warning(f"Missing dependencies: {', '.join(missing)}")
    
    async def extract_from_file(self, file_path: str) -> str:
        """Extract text from a single file"""
        try:
            path = Path(file_path)
            
            if not path.exists():
                return f"Error: File not found - {file_path}"
            
            if not path.is_file():
                return f"Error: Path is not a file - {file_path}"
            
            extension = path.suffix.lower()
            
            if extension not in self.SUPPORTED_EXTENSIONS:
                return f"Error: Unsupported file type - {extension}. Supported: {', '.join(self.SUPPORTED_EXTENSIONS.keys())}"
            
            # Extract text based on file type
            if extension == '.pdf':
                return await self._extract_pdf(path)
            elif extension == '.docx':
                return await self._extract_docx(path)
            elif extension == '.pptx':
                return await self._extract_pptx(path)
            elif extension == '.xlsx':
                return await self._extract_xlsx(path)
            
            return f"Error: Handler not implemented for {extension}"
            
        except Exception as e:
            logger.error(f"Error extracting from {file_path}: {str(e)}")
            return f"Error processing {file_path}: {str(e)}"
    
    async def extract_from_directory(self, directory_path: str, recursive: bool = False) -> str:
        """Extract text from all supported files in a directory"""
        try:
            path = Path(directory_path)
            
            if not path.exists():
                return f"Error: Directory not found - {directory_path}"
            
            if not path.is_dir():
                return f"Error: Path is not a directory - {directory_path}"
            
            # Find all supported files
            files = await self._find_supported_files(path, recursive)
            
            if not files:
                return f"No supported files found in {directory_path}"
            
            # Extract text from each file
            results = []
            for file_path in files:
                logger.info(f"Processing: {file_path}")
                text = await self.extract_from_file(str(file_path))
                
                results.append(f"=== {file_path.name} ===")
                results.append(f"Path: {file_path}")
                results.append(f"Type: {self.SUPPORTED_EXTENSIONS.get(file_path.suffix.lower(), 'Unknown')}")
                results.append("Content:")
                results.append(text)
                results.append("=" * 50)
                results.append("")
            
            return "\n".join(results)
            
        except Exception as e:
            logger.error(f"Error processing directory {directory_path}: {str(e)}")
            return f"Error processing directory {directory_path}: {str(e)}"
    
    async def list_supported_files(self, directory_path: str, recursive: bool = False) -> List[Dict]:
        """List all supported files in a directory with metadata"""
        try:
            path = Path(directory_path)
            
            if not path.exists():
                return [{"error": f"Directory not found - {directory_path}"}]
            
            if not path.is_dir():
                return [{"error": f"Path is not a directory - {directory_path}"}]
            
            files = await self._find_supported_files(path, recursive)
            
            result = []
            for file_path in files:
                stat = file_path.stat()
                result.append({
                    "name": file_path.name,
                    "path": str(file_path),
                    "type": self.SUPPORTED_EXTENSIONS.get(file_path.suffix.lower(), 'Unknown'),
                    "extension": file_path.suffix.lower(),
                    "size_bytes": stat.st_size,
                    "size_mb": round(stat.st_size / (1024 * 1024), 2),
                    "modified": datetime.fromtimestamp(stat.st_mtime).isoformat(),
                    "created": datetime.fromtimestamp(stat.st_ctime).isoformat()
                })
            
            return result
            
        except Exception as e:
            logger.error(f"Error listing files in {directory_path}: {str(e)}")
            return [{"error": f"Error listing files: {str(e)}"}]
    
    async def get_file_info(self, file_path: str) -> Dict:
        """Get detailed information about a file"""
        try:
            path = Path(file_path)
            
            if not path.exists():
                return {"error": f"File not found - {file_path}"}
            
            if not path.is_file():
                return {"error": f"Path is not a file - {file_path}"}
            
            stat = path.stat()
            extension = path.suffix.lower()
            
            info = {
                "name": path.name,
                "path": str(path),
                "extension": extension,
                "type": self.SUPPORTED_EXTENSIONS.get(extension, 'Unsupported'),
                "supported": extension in self.SUPPORTED_EXTENSIONS,
                "size_bytes": stat.st_size,
                "size_mb": round(stat.st_size / (1024 * 1024), 2),
                "modified": datetime.fromtimestamp(stat.st_mtime).isoformat(),
                "created": datetime.fromtimestamp(stat.st_ctime).isoformat()
            }
            
            return info
            
        except Exception as e:
            logger.error(f"Error getting file info for {file_path}: {str(e)}")
            return {"error": f"Error getting file info: {str(e)}"}
    
    async def _find_supported_files(self, path: Path, recursive: bool) -> List[Path]:
        """Find all supported files in a directory"""
        files = []
        
        if recursive:
            for file_path in path.rglob("*"):
                if file_path.is_file() and file_path.suffix.lower() in self.SUPPORTED_EXTENSIONS:
                    files.append(file_path)
        else:
            for file_path in path.iterdir():
                if file_path.is_file() and file_path.suffix.lower() in self.SUPPORTED_EXTENSIONS:
                    files.append(file_path)
        
        return sorted(files)
    
    async def _extract_pdf(self, path: Path) -> str:
        """Extract text from PDF file"""
        if not PDF_AVAILABLE:
            return "Error: PDF library not available. Install with: pip install pdfplumber"
        
        try:
            if PDF_LIBRARY == 'pdfplumber':
                return await self._extract_pdf_pdfplumber(path)
            else:
                return await self._extract_pdf_pypdf2(path)
        except Exception as e:
            return f"Error reading PDF: {str(e)}"
    
    async def _extract_pdf_pdfplumber(self, path: Path) -> str:
        """Extract text using pdfplumber"""
        text_parts = []
        
        with pdfplumber.open(path) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                page_text = page.extract_text()
                if page_text and page_text.strip():
                    text_parts.append(f"--- Page {page_num} ---")
                    text_parts.append(page_text)
                    text_parts.append("")
        
        if not text_parts:
            return "No text content found in PDF"
        
        return "\n".join(text_parts)
    
    async def _extract_pdf_pypdf2(self, path: Path) -> str:
        """Extract text using PyPDF2"""
        text_parts = []
        
        with open(path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            
            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                page_text = page.extract_text()
                if page_text.strip():
                    text_parts.append(f"--- Page {page_num + 1} ---")
                    text_parts.append(page_text)
                    text_parts.append("")
        
        if not text_parts:
            return "No text content found in PDF"
        
        return "\n".join(text_parts)
    
    async def _extract_docx(self, path: Path) -> str:
        """Extract text from DOCX file"""
        if not DOCX_AVAILABLE:
            return "Error: python-docx library not available. Install with: pip install python-docx"
        
        try:
            doc = Document(path)
            text_parts = []
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text)
            
            # Also extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text_parts.append(" | ".join(row_text))
            
            if not text_parts:
                return "No text content found in DOCX"
            
            return "\n".join(text_parts)
            
        except Exception as e:
            return f"Error reading DOCX: {str(e)}"
    
    async def _extract_pptx(self, path: Path) -> str:
        """Extract text from PPTX file"""
        if not PPTX_AVAILABLE:
            return "Error: python-pptx library not available. Install with: pip install python-pptx"
        
        try:
            prs = Presentation(path)
            text_parts = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                slide_text.append(f"--- Slide {slide_num} ---")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                
                if len(slide_text) > 1:  # More than just the slide header
                    text_parts.extend(slide_text)
                    text_parts.append("")
            
            if not text_parts:
                return "No text content found in PPTX"
            
            return "\n".join(text_parts)
            
        except Exception as e:
            return f"Error reading PPTX: {str(e)}"
    
    async def _extract_xlsx(self, path: Path) -> str:
        """Extract text and data from XLSX file"""
        if not XLSX_AVAILABLE:
            return "Error: openpyxl and pandas libraries not available. Install with: pip install openpyxl pandas"
        
        try:
            # Load workbook with openpyxl to get sheet names and structure
            workbook = load_workbook(path, read_only=True)
            text_parts = []
            
            text_parts.append(f"Excel Workbook: {path.name}")
            text_parts.append(f"Number of sheets: {len(workbook.sheetnames)}")
            text_parts.append("="*50)
            text_parts.append("")
            
            # Process each sheet
            for sheet_name in workbook.sheetnames:
                try:
                    # Read sheet with pandas
                    df = pd.read_excel(path, sheet_name=sheet_name, engine='openpyxl')
                    
                    text_parts.append(f"=== Sheet: {sheet_name} ===")
                    text_parts.append(f"Dimensions: {df.shape[0]} rows x {df.shape[1]} columns")
                    
                    if df.empty:
                        text_parts.append("Sheet is empty")
                    else:
                        # Add column names
                        text_parts.append("\nColumns:")
                        text_parts.append(", ".join([str(col) for col in df.columns]))
                        
                        # Add sample data (first 10 rows)
                        text_parts.append("\nSample Data (first 10 rows):")
                        sample_data = df.head(10)
                        
                        # Convert to string representation
                        text_parts.append(sample_data.to_string(index=True, max_cols=None))
                        
                        # Add summary statistics for numeric columns
                        numeric_cols = df.select_dtypes(include=['number']).columns
                        if len(numeric_cols) > 0:
                            text_parts.append("\nNumeric Summary:")
                            summary = df[numeric_cols].describe()
                            text_parts.append(summary.to_string())
                        
                        # Add info about data types
                        text_parts.append("\nData Types:")
                        for col, dtype in df.dtypes.items():
                            non_null_count = df[col].count()
                            total_count = len(df)
                            text_parts.append(f"{col}: {dtype} ({non_null_count}/{total_count} non-null)")
                    
                    text_parts.append("\n" + "="*50)
                    text_parts.append("")
                    
                except Exception as sheet_error:
                    text_parts.append(f"Error reading sheet '{sheet_name}': {str(sheet_error)}")
                    text_parts.append("")
            
            workbook.close()
            
            if len(text_parts) <= 4:  # Only headers, no actual content
                return "No readable content found in XLSX file"
            
            return "\n".join(text_parts)
            
        except Exception as e:
            return f"Error reading XLSX: {str(e)}"
